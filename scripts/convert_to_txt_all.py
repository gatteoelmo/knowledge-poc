import os
import zipfile
from lxml import etree
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

# --- PDF ---
def convert_pdf_to_txt(filepath):
    text = ""
    try:
        with fitz.open(filepath) as pdf:
            for page in pdf:
                text += page.get_text()
    except Exception as e:
        print(f"Errore PDF {filepath}: {e}")
    return text

# --- PPTX ---
def convert_pptx_to_txt(filepath):
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Errore PPTX {filepath}: {e}")
    return text

# --- DOCX ---
def convert_docx_to_txt(filepath):
    text = ""
    try:
        doc = Document(filepath)
        for para in doc.paragraphs:
            if para.text:
                text += para.text + "\n"
    except Exception as e:
        print(f"Errore DOCX {filepath}: {e}")
    return text

# --- KEYNOTE ---
def convert_key_to_txt(filepath):
    text = ""
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            for name in z.namelist():
                if name.endswith(".xml") or name.endswith(".apxl"):
                    try:
                        data = z.read(name)
                        parser = etree.XMLParser(recover=True, encoding='utf-8')
                        root = etree.fromstring(data, parser)
                        for elem in root.iter():
                            if elem.text and elem.text.strip():
                                text += elem.text.strip() + "\n"
                            for attr in elem.attrib.values():
                                if attr.strip():
                                    text += attr.strip() + "\n"
                    except Exception:
                        continue
    except Exception as e:
        print(f"Errore KEY {filepath}: {e}")
    return text

# --- Funzione principale (ricorsiva) ---
def convert_all_in_folder(input_folder, output_root):
    supported_ext = [".pdf", ".pptx", ".key", ".docx"]
    
    for root, dirs, files in os.walk(input_folder):
        # Calcola la sottocartella di output
        rel_path = os.path.relpath(root, input_folder)
        output_folder = os.path.join(output_root, rel_path)
        os.makedirs(output_folder, exist_ok=True)
        
        for filename in files:
            ext = os.path.splitext(filename)[1].lower()
            if ext not in supported_ext:
                continue
            
            filepath = os.path.join(root, filename)
            print(f"Converto: {filepath} ...")
            
            try:
                if ext == ".pdf":
                    text = convert_pdf_to_txt(filepath)
                elif ext == ".pptx":
                    text = convert_pptx_to_txt(filepath)
                elif ext == ".key":
                    text = convert_key_to_txt(filepath)
                elif ext == ".docx":
                    text = convert_docx_to_txt(filepath)
                else:
                    continue

                output_path = os.path.join(output_folder, os.path.splitext(filename)[0] + ".txt")
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)
                
                if not text.strip():
                    print(f"Nessun testo trovato in: {filepath}")
                else:
                    print(f"✅ Salvato: {output_path}")
            
            except Exception as e:
                print(f"❌ Errore con {filepath}: {e}")

# --- Esecuzione ---
if __name__ == "__main__":
    # Struttura cartelle:
    # data/
    # ├── source/      (file originali: PDF, PPTX, KEY, DOCX)
    # └── converted/   (file TXT convertiti)
    
    # Cartella base per i dati
    data_folder = os.path.abspath("./data")
    
    # Cartella con i file da convertire
    input_folder = os.path.join(data_folder, "source")
    
    # Cartella di output per i file convertiti
    output_folder = os.path.join(data_folder, "converted")
    
    # Crea le cartelle se non esistono
    os.makedirs(input_folder, exist_ok=True)
    os.makedirs(output_folder, exist_ok=True)
    
    if not os.path.isdir(input_folder):
        print(f"❌ Cartella non trovata: {input_folder}")
        print(f"💡 Crea la cartella e metti i file da convertire dentro: {input_folder}")
    else:
        print(f"📂 Input:  {input_folder}")
        print(f"📂 Output: {output_folder}\n")
        
        convert_all_in_folder(input_folder, output_folder)
        print(f"\n✅ Tutti i file convertiti sono in: {output_folder}")
